import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Helper for Dark Theme
def set_dark_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 25, 30)

# 1. Custom Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)

txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Smart Route Planning System"
p.font.bold = True
p.font.size = Pt(54)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11.33), Inches(3))
tf2 = txBox2.text_frame
tf2.word_wrap = True

p2 = tf2.paragraphs[0]
p2.text = "Geographical Routing & Emergency Spatial Search"
p2.font.size = Pt(28)
p2.font.color.rgb = RGBColor(100, 200, 255)
p2.alignment = PP_ALIGN.CENTER

p3 = tf2.add_paragraph()
p3.text = "\nSubmitted By: Udayveer Singh Chaudhary (12320106)\nCourse: B.Tech CSE\nSubmitted To: Sachin Garg"
p3.font.size = Pt(24)
p3.font.color.rgb = RGBColor(200, 200, 200)
p3.alignment = PP_ALIGN.CENTER

def add_content_slide(title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_bg(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.73), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = RGBColor(100, 200, 255)
    
    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.73), Inches(5.5))
    tf_body = body_box.text_frame
    tf_body.word_wrap = True
    
    for idx, point in enumerate(bullet_points):
        p = tf_body.add_paragraph() if idx > 0 else tf_body.paragraphs[0]
        if point.startswith("    "):
            p.text = point.strip()
            p.level = 1
            p.font.size = Pt(22)
            p.font.bold = False
        else:
            p.text = "• " + point
            p.level = 0
            p.font.size = Pt(26)
            p.font.bold = True
            
        p.font.color.rgb = RGBColor(230, 230, 230)
        p.space_after = Pt(10)

def add_image_slide(title_text, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_bg(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(100, 200, 255)
    p.alignment = PP_ALIGN.CENTER
    
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.66), Inches(1.2), width=Inches(10))

# 2. Introduction
add_content_slide("1. Introduction & Objectives", [
    "Efficient geographical routing is critical for modern logistics and emergency response.",
    "Problem Statement:",
    "    Commercial maps are heavy and rely on internet APIs.",
    "    Need a dynamic system that adapts to traffic and weather offline.",
    "Objectives:",
    "    Model geographical data (cities/highways) using Graph Theory.",
    "    Implement heuristic pathfinding algorithms from scratch.",
    "    Develop a proximity-based search engine for emergency facilities."
])

# 3. Technology Stack
add_content_slide("2. Technology Stack", [
    "The system was built entirely on native, lightweight technologies:",
    "Programming Language:",
    "    Java (JDK 17) - Chosen for robust Object-Oriented capabilities.",
    "UI Framework:",
    "    Java Swing - Core foundation for the desktop interface.",
    "    FlatLaf Engine - Modern library for a sleek, high-DPI UI.",
    "Architecture Pattern:",
    "    MVC (Model-View-Controller) separating data logic from visual rendering.",
    "Deployment:",
    "    Standalone Executable Java Archive (.jar)"
])

# 4. Architecture (Visual Roadmap)
def add_architecture_visual():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_bg(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.73), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "3. System Architecture Roadmap (MVC)"
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = RGBColor(100, 200, 255)
    
    # Draw Boxes
    # VIEW
    shape_view = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.5), Inches(2.5), Inches(3))
    shape_view.fill.solid()
    shape_view.fill.fore_color.rgb = RGBColor(40, 40, 80)
    tf = shape_view.text_frame
    p = tf.paragraphs[0]
    p.text = "VIEW\n(Frontend UI)\n\n• Java Swing\n• FlatLaf Engine\n• Custom MapPanel"
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER
    
    # CONTROLLER
    shape_ctrl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.4), Inches(2.5), Inches(2.5), Inches(3))
    shape_ctrl.fill.solid()
    shape_ctrl.fill.fore_color.rgb = RGBColor(80, 40, 40)
    tf = shape_ctrl.text_frame
    p = tf.paragraphs[0]
    p.text = "CONTROLLER\n(Logic & DSA)\n\n• RoutingAlgorithms\n• NearbyPlacesFinder\n• CostCalculators"
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER

    # MODEL
    shape_mod = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.3), Inches(2.5), Inches(2.5), Inches(3))
    shape_mod.fill.solid()
    shape_mod.fill.fore_color.rgb = RGBColor(40, 80, 40)
    tf = shape_mod.text_frame
    p = tf.paragraphs[0]
    p.text = "MODEL\n(Database)\n\n• Graph (Adjacency)\n• City Objects\n• Road Objects"
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER

    # Connectors
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.0), Inches(4), Inches(5.4), Inches(4))
    connector.line.width = Pt(5)
    connector.line.color.rgb = RGBColor(255, 255, 255)
    
    connector2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.9), Inches(4), Inches(9.3), Inches(4))
    connector2.line.width = Pt(5)
    connector2.line.color.rgb = RGBColor(255, 255, 255)

add_architecture_visual()

# 5. File Structure (Visual Roadmap)
def add_filestructure_visual():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_bg(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.73), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "4. File & Package Roadmap"
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = RGBColor(100, 200, 255)
    
    # Root
    shape_root = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.0), Inches(2.8), Inches(1))
    shape_root.fill.solid()
    shape_root.fill.fore_color.rgb = RGBColor(60, 60, 60)
    tf = shape_root.text_frame
    p = tf.paragraphs[0]
    p.text = "com.routeplanner"
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER

    packages = [
        ("model", "Entity Data Classes (City, Road, Graph)", 1.5),
        ("dao", "Database Access & Storage", 2.6),
        ("dsa", "Core Algorithm Engines", 3.7),
        ("manager", "Business Logic & Computation", 4.8),
        ("ui", "Graphical Interface Components", 5.9)
    ]
    
    for i, (pkg, desc, y_pos) in enumerate(packages):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(y_pos), Inches(8.5), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(40, 80, 120)
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{pkg}   ➔   {desc}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        
        # Connector from root (X=3.3) to Box (X=4.2)
        conn = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, Inches(3.3), Inches(4.5), Inches(4.2), Inches(y_pos + 0.45))
        conn.line.width = Pt(3)
        conn.line.color.rgb = RGBColor(200, 200, 200)

add_filestructure_visual()

# 6. Core Data Structures
add_content_slide("5. Core Data Structures (DSA)", [
    "Heavily optimized data structures form the backbone of the system:",
    "Adjacency List (Graphs):",
    "    Stores the network of cities and roads (edges).",
    "    Ensures memory scales linearly O(V + E) rather than exponentially.",
    "Priority Queues (Min-Heap):",
    "    Used in Dijkstra/A* to extract the closest nodes in O(log V) time.",
    "Priority Queues (Max-Heap):",
    "    Maintains the K-nearest emergency facilities during spatial searches.",
    "HashMaps:",
    "    Caches city names to coordinate objects for instant O(1) lookups."
])

# 7. Core Algorithms
add_content_slide("6. Core Algorithms Implemented", [
    "Multiple algorithms were implemented and compared:",
    "Dijkstra's Algorithm:",
    "    Calculates the absolute shortest path across all edges.",
    "A* (A-Star) Search:",
    "    Utilizes the Haversine formula as a geographic heuristic to boost speed.",
    "Breadth-First Search (BFS):",
    "    Ignores physical distance to find routes with the fewest stops.",
    "Floyd-Warshall:",
    "    Dynamic programming for all-pairs shortest paths."
])

# 8. Haversine Formula
add_content_slide("7. The Haversine Formula", [
    "Crucial mathematical logic used to account for Earth's curvature:",
    "Purpose:",
    "    Calculates the 'great-circle' geographic distance between two points.",
    "Usage in System:",
    "    Acts as the heuristic 'h(n)' in the A* Search algorithm.",
    "    Validates whether an emergency facility is strictly within the 50km radius.",
    "Formula Concept:",
    "    a = sin²(Δφ/2) + cos φ1 ⋅ cos φ2 ⋅ sin²(Δλ/2)",
    "    c = 2 ⋅ atan2( √a, √(1−a) )",
    "    d = R ⋅ c (Where R = Earth Radius 6371km)"
])

# Add Image Slides (Using uploaded images)
add_image_slide("Route Planner: Shortest Path Demonstration", "/Users/mac/.gemini/antigravity/brain/90e08bfa-22d0-430e-bc98-bfbf1932a50e/.user_uploaded/media_1786020561530.png")
add_image_slide("Spatial Search: Nearby Emergency Locator", "/Users/mac/.gemini/antigravity/brain/90e08bfa-22d0-430e-bc98-bfbf1932a50e/.user_uploaded/media_1786020583608.png")
add_image_slide("Database Management: Administrator Panel", "/Users/mac/.gemini/antigravity/brain/90e08bfa-22d0-430e-bc98-bfbf1932a50e/.user_uploaded/media_1786020612010.png")

# 12. Future Scope
add_content_slide("8. Future Scope", [
    "The modular architecture allows for massive future scalability:",
    "Live API Integration:",
    "    Fetch real-time traffic data from Google Maps API or OpenWeather API.",
    "Database Persistence:",
    "    Migrate in-memory HashMaps to an AWS RDS (MySQL) cloud database.",
    "Turn-by-Turn Telemetry:",
    "    Implement vector calculations to generate 'Turn Left / Turn Right' instructions.",
    "Mobile Porting:",
    "    Migrate core Java DSA logic to Kotlin for Android deployment."
])

# 13. Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)
txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(3))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Thank You!"
p.font.bold = True
p.font.size = Pt(72)
p.font.color.rgb = RGBColor(100, 200, 255)
p.alignment = PP_ALIGN.CENTER

prs.save('/Users/mac/Desktop/SmartRoutePlanner/Project_Presentation_V4.pptx')
