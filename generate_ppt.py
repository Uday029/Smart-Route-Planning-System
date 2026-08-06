import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()

def add_slide(title_text, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = bullet_points[0]
    
    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        if point.startswith("    "):
            p.text = point.strip()
            p.level = 1
        else:
            p.text = point
            p.level = 0
            
# 1. Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Smart Route Planning System"
subtitle.text = "Geographical Routing & Emergency Spatial Search\n\nSubmitted By: Udayveer Singh Chaudhary (12320106)\nCourse: B.Tech CSE\nSubmitted To: Sachin Garg"

# 2. Introduction
add_slide("1. Introduction", [
    "Efficient geographical routing is critical for modern logistics and emergency response.",
    "Problem Statement:",
    "    • Commercial maps are heavy and rely on internet APIs.",
    "    • Need a dynamic system that adapts to traffic and weather offline.",
    "Objectives:",
    "    • Model geographical data (cities/highways) using Graph Theory.",
    "    • Implement heuristic pathfinding algorithms from scratch.",
    "    • Develop a proximity-based search engine for emergency facilities."
])

# 3. Technology Stack
add_slide("2. Technology Stack", [
    "The system was built entirely on native, lightweight technologies:",
    "Programming Language:",
    "    • Java (JDK 17) - Chosen for robust Object-Oriented capabilities.",
    "UI Framework:",
    "    • Java Swing - Core foundation for the desktop interface.",
    "    • FlatLaf Engine - Modern 3rd-party library for a sleek, high-DPI UI.",
    "Architecture Pattern:",
    "    • MVC (Model-View-Controller) separating data logic from visual rendering.",
    "Deployment:",
    "    • Standalone Executable Java Archive (.jar)"
])

# 4. Core Data Structures
add_slide("3. Core Data Structures (DSA)", [
    "Heavily optimized data structures form the backbone of the system:",
    "Adjacency List (Graphs):",
    "    • Stores the network of cities and roads.",
    "    • Ensures memory scales linearly O(V + E) rather than exponentially.",
    "Priority Queues (Min-Heap):",
    "    • Used in Dijkstra/A* to extract the closest nodes in O(log V) time.",
    "Priority Queues (Max-Heap):",
    "    • Maintains the K-nearest emergency facilities during spatial searches.",
    "HashMaps:",
    "    • Caches city names to coordinate objects for instant O(1) lookups."
])

# 5. Core Algorithms
add_slide("4. Core Algorithms Implemented", [
    "Multiple algorithms were implemented and compared:",
    "Dijkstra's Algorithm:",
    "    • Calculates the absolute shortest path across all edges.",
    "A* (A-Star) Search:",
    "    • Utilizes the Haversine formula as a geographic heuristic to boost speed.",
    "Breadth-First Search (BFS):",
    "    • Ignores physical distance to find routes with the fewest stops.",
    "Floyd-Warshall:",
    "    • Dynamic programming for all-pairs shortest paths."
])

# 6. Haversine Formula
add_slide("5. The Haversine Formula", [
    "Crucial mathematical logic used to account for Earth's curvature:",
    "Purpose:",
    "    • Calculates the 'great-circle' geographic distance between two points.",
    "Usage in System:",
    "    • Acts as the heuristic 'h(n)' in the A* Search algorithm.",
    "    • Validates whether an emergency facility is strictly within the 50km radius.",
    "Formula Concept:",
    "    • a = sin²(Δφ/2) + cos φ1 ⋅ cos φ2 ⋅ sin²(Δλ/2)",
    "    • c = 2 ⋅ atan2( √a, √(1−a) )",
    "    • d = R ⋅ c (Where R = Earth Radius 6371km)"
])

# 7. Module: Route Planner
add_slide("6. Module: Route Planner", [
    "The primary navigation interface of the software:",
    "Features:",
    "    • Dynamic Edge Modifiers: Simulates Live Traffic (delays) or Bad Weather (roadblocks) and recalculates the path instantly.",
    "    • Cost Estimation Engine: Calculates estimated fuel consumption (Liters) and tolls.",
    "    • Custom 2D Graphics Engine: Normalizes Lat/Lon coordinates to screen pixels.",
    "    • Auto-Zoom: Mathematically calculates a bounding box to frame the final route perfectly on the map."
])

# 8. Module: Nearby & Emergency
add_slide("7. Module: Nearby & Emergency", [
    "A localized spatial search engine for critical infrastructure:",
    "Features:",
    "    • Queries the graph for Hospitals, Police Stations, and Petrol Pumps.",
    "    • Hard limit of 50 km search radius.",
    "    • Max-Heap implementation discards distant locations to find the K-nearest.",
    "    • Numbered visual map markers map exactly to the printed text list to prevent UI overlap."
])

# 9. Module: Admin Panel
add_slide("8. Module: Admin Data Injector", [
    "Allows real-time expansion of the graph database without code changes:",
    "Add Cities:",
    "    • Automatically assigns IDs and plots new coordinate points on the map.",
    "Add Roads:",
    "    • Forges connecting edges between cities with specific speed limits and distances.",
    "Add Places:",
    "    • Spawns localized facilities (e.g., new ATMs) that are instantly searchable by the Nearby engine."
])

# 10. Future Scope
add_slide("9. Future Scope", [
    "The modular architecture allows for massive future scalability:",
    "Live API Integration:",
    "    • Fetch real-time traffic data from Google Maps API or OpenWeather API.",
    "Database Persistence:",
    "    • Migrate in-memory HashMaps to an AWS RDS (MySQL) cloud database.",
    "Turn-by-Turn Telemetry:",
    "    • Implement vector calculations to generate 'Turn Left / Turn Right' instructions.",
    "Mobile Porting:",
    "    • Migrate core Java DSA logic to Kotlin for Android deployment."
])

# 11. Conclusion
add_slide("10. Conclusion", [
    "Summary of project achievements:",
    "    • Successfully bridged theoretical computer science with practical software engineering.",
    "    • Proved the massive efficiency gains of optimized Heaps and Adjacency Lists.",
    "    • Built a fully-featured, deployable desktop application.",
    "    • The system is highly scalable, interactive, and parallels professional logistics software.",
    "",
    "Thank You!"
])

prs.save('/Users/mac/Desktop/SmartRoutePlanner/Project_Presentation.pptx')
